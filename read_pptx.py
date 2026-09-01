from pptx import Presentation

prs = Presentation('OceanEmbed_SIH_Submission.pptx.pptx')
for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            print(shape.text)
